from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
from sambanova import generate_summary
from supabase_client import save_note
from pptx import Presentation
from PyPDF2 import PdfReader
import io

app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allows all origins
    allow_credentials=True,
    allow_methods=["*"],  # Allows all methods
    allow_headers=["*"],  # Allows all headers
)

class NoteRequest(BaseModel):
    title: str
    content: str
    user_id: str  # UUID from auth.users
    folder_id: str  # UUID from folders table
    source_type: str  # 'file', 'youtube', 'audio', 'text', 'live'
    source_url: Optional[str] = None
    language: Optional[str] = None

def extract_text_from_pptx(file_content: bytes) -> str:
    """Extract text from a PowerPoint file."""
    try:
        # Create a presentation object from the file content
        prs = Presentation(io.BytesIO(file_content))
        
        # Extract text from each slide
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        print(text_content)
        return "\n".join(text_content)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error processing PowerPoint file: {str(e)}")

def extract_text_from_pdf(file_content: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        # Create a PDF reader object from the file content
        pdf_reader = PdfReader(io.BytesIO(file_content))
        
        # Extract text from each page
        text_content = []
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:  # Only add non-empty text
                text_content.append(text)
        
        if not text_content:
            raise ValueError("No text could be extracted from the PDF")
            
        return "\n".join(text_content)
    except Exception as e:
        print(f"Error in PDF extraction: {e}")
        raise HTTPException(status_code=400, detail=f"Error processing PDF file: {str(e)}")

@app.post("/process-ppt")
async def process_ppt(
    file: UploadFile = File(...),
    user_id: str = Form(...),
    folder_id: str = Form(...),
    source_type: str = Form("file")
):
    try:
        # Read the file content
        file_content = await file.read()
        
        # Extract text from PowerPoint
        text_content = extract_text_from_pptx(file_content)
        
        # Generate summary
        summary = generate_summary(text_content)
        
        # Prepare note data
        note_data = {
            "title": summary["title"],
            "content": text_content,
            "user_id": user_id,
            "folder_id": folder_id,
            "source_type": source_type,
            "source_url": None,
            "language": None,
            "summary": summary["summary"],
            "quiz_questions": None,
        }
        
        # Save note to Supabase
        saved_note = save_note(note_data)
        if not saved_note:
            raise HTTPException(status_code=500, detail="Failed to save note to database")
            
        return {
            "status": "success",
            "note": saved_note,
            "summary": summary
        }

    except Exception as e:
        print(f"Error processing PowerPoint: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/process-note")
async def process_note(payload: NoteRequest):
    try:
        print("Processing note with content:", payload.content)
        summary = generate_summary(payload.content)
        
        # Validate the summary response
        if not isinstance(summary, dict):
            raise HTTPException(status_code=500, detail="Invalid summary response format")
            
        if "title" not in summary:
            raise HTTPException(status_code=500, detail="Summary missing title")

        # Generate embedding for the content
        # embedding = generate_embedding(payload.content)
        
        # Convert summary to text format
        
        
        note_data = {
            "title": summary["title"],
            "content": payload.content,
            "user_id": payload.user_id,
            "folder_id": payload.folder_id,
            "source_type": payload.source_type,
            "source_url": payload.source_url,
            "language": payload.language,
            "summary": summary["summary"],  # Now a plain text string #
            "quiz_questions": None,  # Can be updated later
        }
        
        print("Processed note data:", note_data)
        
        # Save note to Supabase
        saved_note = save_note(note_data)
        if not saved_note:
            raise HTTPException(status_code=500, detail="Failed to save note to database")
            
        return {"status": "success", "note": saved_note}

    except Exception as e:
        print(f"Error processing note: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/process-pdf")
async def process_pdf(
    content: UploadFile = File(...),
    user_id: str = Form(...),
    folder_id: str = Form(...),
    source_type: str = Form("file")
):
    try:
        # Read the file content
        file_content = await content.read()
        
        # Extract text from PDF
        text_content = extract_text_from_pdf(file_content)
        
        # Generate summary
        summary = generate_summary(text_content)
        
        # Prepare note data
        note_data = {
            "title": summary["title"],
            "content": text_content,
            "user_id": user_id,
            "folder_id": folder_id,
            "source_type": source_type,
            "source_url": None,
            "language": None,
            "summary": summary["summary"],
            "quiz_questions": None,
        }
        
        # Save note to Supabase
        saved_note = save_note(note_data)
        if not saved_note:
            raise HTTPException(status_code=500, detail="Failed to save note to database")
        
        return {"status": "success", "note": saved_note}

    except Exception as e:
        print(f"Error processing PDF: {e}")
        raise HTTPException(status_code=500, detail=str(e))

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8001)